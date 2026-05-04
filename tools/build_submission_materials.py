from pathlib import Path

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DInches
from docx.shared import Pt as DPt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("submission_materials")
OUT.mkdir(exist_ok=True)


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    bg = RGBColor(246, 248, 250)
    ink = RGBColor(18, 31, 43)
    muted = RGBColor(91, 105, 120)
    blue = RGBColor(38, 103, 255)
    green = RGBColor(22, 163, 74)
    red = RGBColor(220, 38, 38)
    yellow = RGBColor(234, 179, 8)
    panel = RGBColor(255, 255, 255)
    line = RGBColor(210, 218, 226)
    dark = RGBColor(13, 27, 42)

    def set_bg(slide, color=bg):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def text(slide, x, y, w, h, value, size=24, bold=False, color=ink, align=None):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = value
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return box

    def title(slide, value, subtitle=None, idx=None):
        text(slide, 0.65, 0.42, 9.6, 0.55, value, 28, True, ink)
        if subtitle:
            text(slide, 0.67, 0.98, 10.8, 0.42, subtitle, 13, False, muted)
        if idx is not None:
            text(slide, 12.2, 0.5, 0.45, 0.25, f"{idx:02d}", 10, False, muted, PP_ALIGN.RIGHT)

    def card(slide, x, y, w, h, heading, body, accent=blue):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = panel
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
        shape.adjustments[0] = 0.08
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        text(slide, x + 0.25, y + 0.18, w - 0.45, 0.28, heading, 14, True, ink)
        if body:
            text(slide, x + 0.25, y + 0.58, w - 0.45, h - 0.72, body, 11.5, False, muted)

    def placeholder(slide, x, y, w, h, label):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(236, 240, 244)
        shape.line.color.rgb = RGBColor(175, 184, 193)
        shape.line.width = Pt(1.2)
        shape.adjustments[0] = 0.04
        text(slide, x + 0.25, y + h / 2 - 0.18, w - 0.5, 0.4, label, 13, True, muted, PP_ALIGN.CENTER)

    def bullets(slide, x, y, w, h, items, size=15):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            run = p.add_run()
            run.text = "• " + item
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.color.rgb = ink

    slide = prs.slides.add_slide(blank)
    set_bg(slide, dark)
    text(slide, 0.85, 0.75, 8.2, 0.55, "Smart-Home-Guard", 18, True, RGBColor(139, 210, 255))
    text(slide, 0.82, 1.65, 8.6, 1.4, "端侧智能居家安全告警系统", 44, True, RGBColor(255, 255, 255))
    text(slide, 0.86, 3.08, 7.8, 0.55, "Aurora 画面配置 · 端侧检测 · 危险区告警 · 60 秒可验收日志", 18, False, RGBColor(190, 207, 224))
    for x, y, w, h, color in [(9.2, 1.0, 3.1, 2.1, blue), (8.85, 3.55, 2.4, 1.55, yellow), (10.5, 4.65, 1.4, 0.9, red)]:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        box.fill.background()
        box.line.color.rgb = color
        box.line.width = Pt(5)
    text(slide, 0.9, 6.55, 5.8, 0.35, "初赛提交材料草稿 | 需补充现场截图与最终实测日志", 12, False, RGBColor(155, 175, 195))

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "重要内容预览", "把评审最关心的“能跑、可验、稳定、可复现”提前亮出来", 2)
    card(slide, 0.75, 1.65, 3.8, 1.45, "功能闭环", "Aurora 取景 -> PC 画危险区 -> 板端检测 -> OSD 显示 -> GPIO 声光报警。", blue)
    card(slide, 4.8, 1.65, 3.8, 1.45, "性能可量化", "60 秒自动输出 FPS、R 值、P95 延迟、评分估算和异常计数。", green)
    card(slide, 8.85, 1.65, 3.8, 1.45, "现场鲁棒", "摄像头、数据、推理、资源异常均有日志、降级和恢复策略。", red)
    placeholder(slide, 0.75, 3.55, 5.8, 2.6, "待补：Aurora 调试画面截图")
    placeholder(slide, 6.85, 3.55, 5.8, 2.6, "待补：最终串口 [CHECK] 日志截图")

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "设计目标", "围绕居家危险区域入侵告警，优先保证现场可复现和评审可验证", 3)
    bullets(slide, 0.9, 1.65, 6.0, 3.9, [
        "检测目标：person、dog、cat 三类进入危险区触发告警",
        "部署目标：端侧独立运行，PC 仅用于配置危险区",
        "可视化目标：OSD 显示危险区、目标框和 ALERT 状态",
        "验收目标：稳定运行 >= 60s，输出 FPS/P95 延迟/异常计数",
        "鲁棒目标：黑暗、强光、资源缺失和链路异常下可降级运行",
    ], 17)
    card(slide, 7.25, 1.75, 5.0, 1.1, "赛题评分映射", "功能目标完成度 40 + 应用性能 20 + 创新性 20 + 协作与材料 20", blue)
    card(slide, 7.25, 3.15, 5.0, 1.1, "取舍原则", "不追求复杂 UI，优先降低端侧热路径 I/O 与推理延迟。", green)
    card(slide, 7.25, 4.55, 5.0, 1.1, "边界说明", "多边形判断是真实多边形；OSD 展示为外接矩形。", yellow)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "系统架构", "板端完成实时检测与告警，PC 端只参与配置与调试", 4)
    labels = [("Camera / Aurora", 0.8), ("A1 图像管线", 2.9), ("AI 推理", 5.0), ("区域判断", 7.0), ("OSD/GPIO", 9.0), ("串口日志", 11.0)]
    for label, x in labels:
        card(slide, x, 2.3, 1.55, 0.95, label, "", blue)
    for x in [2.35, 4.45, 6.45, 8.45, 10.45]:
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(2.78), Inches(x + 0.45), Inches(2.78))
        connector.line.color.rgb = muted
        connector.line.width = Pt(2)
    card(slide, 0.95, 4.45, 5.1, 1.35, "PC 上位机", "串口 SNAPSHOT 获取预览，鼠标绘制区域，发送 ZONE/START。", green)
    card(slide, 7.15, 4.45, 5.1, 1.35, "板端运行态", "正式检测阶段压缩快照频率，验收日志直接输出到串口。", red)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "实现思路", "把“检测”“判断”“显示”“告警”“验收”拆成可观察模块", 5)
    modules = [
        (0.8, 1.55, "检测模块", "Head6 输出 -> DFL 解码 -> stride 恢复 -> NMS -> 类别过滤。", blue),
        (4.8, 1.55, "危险区模块", "PC 坐标保持在 1440x1080 crop 空间，板端同坐标判断。", green),
        (8.8, 1.55, "可视化模块", "OSD 分层绘制检测框、危险区和 ALERT 位图。", yellow),
        (0.8, 3.25, "报警模块", "进入危险区后 GPIO8/GPIO10 控制 LED 与蜂鸣器。", red),
        (4.8, 3.25, "环境策略", "采样 UYVY Y 通道，低照/强光动态调整阈值。", blue),
        (8.8, 3.25, "验收日志", "60s 汇总 runtime、avg_app、R、p95_T、异常计数。", green),
    ]
    for x, y, h, b, c in modules:
        card(slide, x, y, 3.8, 1.25, h, b, c)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "模型与端侧适配", "从通用检测基线到 SmartSens M1 部署模型", 6)
    text(slide, 0.85, 1.6, 5.8, 0.4, "正式模型文件", 16, True, ink)
    text(slide, 0.9, 2.1, 5.8, 0.4, "smart_guard_coco_256.m1model", 24, True, blue)
    card(slide, 0.85, 2.9, 5.4, 1.4, "端侧优化", "输入尺寸从 640 降为 256，减少预处理和推理负载；检测后处理保留 DFL/stride/NMS 的完整链路。", green)
    card(slide, 6.85, 1.65, 5.7, 1.2, "真实声明", "通用 COCO 检测权重作为基线，结合 Aurora 现场画面进行阈值、尺寸和鲁棒性调优。", yellow)
    card(slide, 6.85, 3.15, 5.7, 1.2, "不夸大边界", "当前材料不声称已完成 Aurora 自采样再训练；如需满足强制要求，应补采样、标注和微调证据。", red)
    placeholder(slide, 0.85, 4.8, 11.7, 1.0, "待补：模型转换平台截图或最终模型文件列表截图")

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "危险区判断与 OSD 显示边界", "判断追求准确，显示追求稳定，二者在工程中明确分离", 7)
    card(slide, 0.85, 1.55, 5.2, 1.35, "真实多边形判断", "多边形点位来自 PC 上位机；板端用 point-in-polygon 判断目标中心是否进入区域。", green)
    card(slide, 0.85, 3.15, 5.2, 1.35, "外接矩形显示", "OSD 当前显示多边形的 bounding box，避免多段线绘制造成 buffer 与性能波动。", yellow)
    card(slide, 0.85, 4.75, 5.2, 1.0, "日志可见", "串口输出：[ZONE] judgement=polygon display=bbox points=N", blue)
    placeholder(slide, 6.6, 1.55, 5.9, 4.2, "待补：危险区画框截图（标注真实点集与外接矩形）")

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "PC 上位机与串口配置", "无网络时仍可完成危险区配置和启动", 8)
    for i, (num, head, body) in enumerate([("1", "SNAPSHOT", "板端发送 Aurora 预览图"), ("2", "DRAW", "PC 鼠标绘制多边形"), ("3", "ZONE", "发送区域 JSON"), ("4", "START", "进入正式检测")]):
        card(slide, 0.85 + i * 3.1, 1.75, 2.55, 1.35, f"{num}. {head}", body, [blue, green, yellow, red][i])
    placeholder(slide, 0.85, 3.75, 5.7, 2.1, "待补：PC 上位机窗口截图")
    placeholder(slide, 6.85, 3.75, 5.7, 2.1, "待补：FT232RL 接线/实物图")

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "应用性能评估", "最终分数需要用 256 模型重新跑 60 秒日志确认", 9)
    rows = [
        ("指标", "日志字段", "当前取证要求"),
        ("实时性", "[FPS] app/sensor/R/score_est", "保存 60s 内多条 FPS 输出"),
        ("端到端延迟", "[LAT] p95/p95_T/score_est", "保存 P95 与 45fps 帧周期换算"),
        ("稳定性", "[CHECK][SUMMARY] stable", "runtime >= 60s 且无连续失效"),
        ("异常计数", "[CHECK][COUNTS] recoveries/failures", "资源异常为 0 或可解释降级"),
    ]
    table = slide.shapes.add_table(len(rows), 3, Inches(0.85), Inches(1.65), Inches(11.7), Inches(3.3)).table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(229, 236, 244) if r == 0 else RGBColor(255, 255, 255)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(12 if r else 13)
                    run.font.bold = r == 0
                    run.font.color.rgb = ink
    text(slide, 0.9, 5.35, 11.5, 0.6, "待补：最终 256 模型实测数据。不要沿用旧版 640/60fps 数据。", 15, True, red)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "鲁棒性与异常处理", "评分表要求的摄像头、数据、推理、资源异常均有对应策略", 10)
    robust = [
        (0.85, 1.55, "摄像头异常", "连续取帧失败后重启 pipeline，并输出 [CAM][ALARM]。", red),
        (4.85, 1.55, "数据异常", "sensor tensor 空指针或异常时跳过并尝试恢复。", yellow),
        (8.85, 1.55, "推理异常", "预处理/推理失败计数，连续失败时降级跳过 OSD。", blue),
        (0.85, 3.25, "资源异常", "模型、LUT、bitmap、GPIO、UART 缺失均有 warn/degrade。", green),
        (4.85, 3.25, "黑暗/强光", "avg_luma 触发 LOW_LIGHT/BRIGHT，动态调整 conf。", blue),
        (8.85, 3.25, "启动恢复", "run.sh 直接输出串口日志，异常退出后 2s 自动重启。", green),
    ]
    for x, y, h, b, c in robust:
        card(slide, x, y, 3.65, 1.35, h, b, c)

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "问题闭环与优化路径", "从现象、定位、修复到验证，形成可答辩的问题链", 11)
    bullets(slide, 0.9, 1.5, 6.1, 4.8, [
        "人体框不完整：定位到 YOLOv8 DFL 解码缺失 stride，恢复 ltrb*stride_f。",
        "危险区误判：统一 crop 坐标判断，OSD 显示前再加 x offset。",
        "日志影响性能：关闭热路径文件写入，快照编码移到后台线程。",
        "端侧 FPS 偏低：输入尺寸切到 256，降低模型和后处理负载。",
        "现场光照变化：用 avg_luma 轻量判断低照/强光并动态调阈值。",
    ], 15)
    placeholder(slide, 7.35, 1.65, 5.0, 3.9, "待补：优化前后对比截图或 Git 提交记录截图")

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    title(slide, "演示与提交材料计划", "把能证明作品的证据全部沉淀到 PPT、技术文档和压缩包里", 12)
    card(slide, 0.85, 1.55, 3.65, 1.55, "PPT 介绍视频", "录制本 PPT 讲解，控制在 5-8 分钟，覆盖目标、实现、性能、问题闭环。", blue)
    card(slide, 4.85, 1.55, 3.65, 1.55, "作品演示视频", "横屏 1080P，录 Aurora 画面、串口日志、危险区告警和声光反馈。", green)
    card(slide, 8.85, 1.55, 3.65, 1.55, "技术数据包", "源码、模型、脚本、配置、可执行文件、README、最终日志截图。", yellow)
    placeholder(slide, 0.85, 3.75, 11.7, 1.65, "待补：最终材料目录截图 + 演示视频封面帧")
    text(slide, 0.9, 6.35, 11.0, 0.35, "下一步：封版提交前，运行 git status，确认模型文件、材料文件和 zone_config 是否按需入仓。", 12, False, muted)

    path = OUT / "Smart-Home-Guard_intro_draft.pptx"
    prs.save(path)
    return path


def build_docx() -> Path:
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = DInches(0.65)
    sec.bottom_margin = DInches(0.65)
    sec.left_margin = DInches(0.75)
    sec.right_margin = DInches(0.75)
    styles = doc.styles
    styles["Normal"].font.name = "Microsoft YaHei"
    styles["Normal"]._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    styles["Normal"].font.size = DPt(10.5)
    for st in ["Heading 1", "Heading 2", "Heading 3"]:
        styles[st].font.name = "Microsoft YaHei"
        styles[st]._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")

    def shade(cell, fill):
        tc_pr = cell._tc.get_or_add_tcPr()
        shd = OxmlElement("w:shd")
        shd.set(qn("w:fill"), fill)
        tc_pr.append(shd)

    def table(headers, rows):
        tbl = doc.add_table(rows=1, cols=len(headers))
        tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
        tbl.style = "Table Grid"
        for i, head in enumerate(headers):
            cell = tbl.rows[0].cells[i]
            cell.text = head
            shade(cell, "E5ECF4")
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            for p in cell.paragraphs:
                for r in p.runs:
                    r.bold = True
                    r.font.name = "Microsoft YaHei"
                    r._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
        for row in rows:
            cells = tbl.add_row().cells
            for i, value in enumerate(row):
                cells[i].text = str(value)
                cells[i].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
                for p in cells[i].paragraphs:
                    for r in p.runs:
                        r.font.name = "Microsoft YaHei"
                        r._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
        doc.add_paragraph("")

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Smart-Home-Guard\n技术报告")
    run.bold = True
    run.font.size = DPt(24)
    run.font.name = "Microsoft YaHei"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    doc.add_paragraph("端侧智能居家安全告警系统 | 初赛提交材料草稿").alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("说明：本文档已包含核心技术路线和材料结构，实物图、Aurora 调试截图、最终 256 模型 60 秒验收日志需由现场测试后补充。").alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("1. 设计目标", 1)
    doc.add_paragraph("Smart-Home-Guard 面向居家安全场景，在 SmartSens FlyingChip A1 平台完成端侧目标检测、危险区判断、OSD 显示与 GPIO 声光报警。设计重点是现场可运行、结果可验证、参数可配置、异常可恢复。")
    table(["目标项", "实现说明", "评分关联"], [
        ["核心功能", "person/dog/cat 进入危险区触发告警", "功能目标完成度"],
        ["可视化", "检测框、危险区、ALERT 位图分图层 OSD 输出", "功能覆盖与可验证性"],
        ["配置性", "PC 上位机通过串口快照绘制危险区并下发", "复现性与工程深度"],
        ["性能", "60 秒输出 FPS、R、P95 延迟和评分估算", "应用性能评估"],
        ["鲁棒性", "摄像头、数据、推理、资源异常均有日志和降级策略", "异常处理"],
    ])

    doc.add_heading("2. 系统架构", 1)
    doc.add_paragraph("系统由板端主程序、AI 检测模型、OSD/GPIO 外设控制、UART 配置通道和 PC 上位机组成。PC 端只负责配置危险区，正式检测在板端独立完成。")
    table(["模块", "文件/目录", "职责"], [
        ["板端主程序", "ssne_ai_yolo_coco/demo_yolo_coco.cpp", "取流、检测、区域判断、日志和异常处理"],
        ["检测器", "ssne_ai_yolo_coco/src/coco_detector.cpp", "Head6 输出解析、DFL 解码、NMS"],
        ["OSD 工具", "ssne_ai_yolo_coco/src/utils.cpp", "检测框、危险区、ALERT 位图显示"],
        ["报警控制", "ssne_ai_yolo_coco/src/gpio_alarm_controller.cpp", "LED/蜂鸣器 GPIO 控制"],
        ["上位机", "pc_controller/zone_controller.py", "串口快照、危险区绘制和下发"],
    ])

    sections = [
        ("3. 实现思路", "检测结果在 1440x1080 crop 坐标中完成跟踪、类别过滤和危险区判断；只有在 OSD 显示和串口日志展示时才转换到 1920x1080 全图坐标。该策略避免了坐标系混用导致的误判。"),
        ("3.1 检测后处理", "模型输出 P3/P4/P5 三个尺度的 box/class head。板端执行 sigmoid、DFL softmax 期望、stride 恢复、坐标缩放和 NMS。历史问题中，人体框不完整的根因是 DFL ltrb 距离缺失 stride 乘法，当前已恢复 ltrb * stride_f。"),
        ("3.2 危险区判断与显示边界", "矩形危险区按中心点是否落入矩形判断；多边形危险区按 point-in-polygon 判断，使用真实多边形点集。当前正式演示中，多边形在 OSD 上显示为外接矩形，这是为了控制 OSD DMA buffer 和现场稳定性；该显示不改变告警判断结果。串口会输出 [ZONE] judgement=polygon display=bbox points=N 作为显式说明。"),
        ("3.3 PC 配置流程", "板端启动后等待串口命令。PC 端发送 SNAPSHOT 获取 Aurora 预览图，绘制危险区后发送 ZONE <json>，最后发送 START 进入正式检测。该方式不依赖网络，适合现场演示。"),
    ]
    for heading, body in sections:
        doc.add_heading(heading, 1 if heading.startswith("3.") and heading.count(".") == 1 else 2)
        doc.add_paragraph(body)

    doc.add_heading("4. 功能实现情况", 1)
    table(["功能", "当前状态", "验证方式"], [
        ["目标检测", "已实现 person/dog/cat 检测与类别过滤", "Aurora 画面 + 串口 [DET]"],
        ["危险区告警", "已实现矩形/多边形区域判断", "[ALARM] 与 OSD 红框"],
        ["OSD 显示", "检测框、危险区、ALERT 图标", "Aurora 预览截图"],
        ["GPIO 报警", "LED/蜂鸣器联动", "实物视频"],
        ["串口配置", "SNAPSHOT/ZONE/START", "PC 上位机截图"],
        ["验收日志", "60 秒 SUMMARY/COUNTS", "串口日志截图"],
    ])

    doc.add_heading("5. 性能测试数据", 1)
    doc.add_paragraph("最终性能必须使用 256 输入模型重新测试。建议记录至少 60 秒，保存完整串口日志和截图。")
    table(["指标", "日志字段", "填写位置"], [
        ["平均应用 FPS", "[CHECK][SUMMARY] avg_app", "待补：最终 256 实测"],
        ["实时性 R", "R=avg_app/45", "待补：最终 256 实测"],
        ["P95 延迟", "p95 / p95_T", "待补：最终 256 实测"],
        ["稳定性", "stable=PASS/PASS_WITH_RECOVERY", "待补：最终 256 实测"],
        ["异常计数", "cam/data/infer/resource warnings", "待补：最终 256 实测"],
    ])

    doc.add_heading("6. 鲁棒性和异常处理", 1)
    table(["异常类别", "检测方式", "处理策略"], [
        ["摄像头异常", "GetImage 失败计数", "连续失败后重启 pipeline，输出 [CAM][ALARM]"],
        ["数据异常", "sensor tensor 空指针检查", "连续异常后重启 pipeline，输出 [DATA][ALARM]"],
        ["推理异常", "预处理/推理返回码", "计数并跳过本轮 OSD，输出 [INFER][ALARM]"],
        ["资源异常", "模型/LUT/bitmap/GPIO/UART 文件或初始化检查", "检测可降级，输出 [RESOURCE][WARN/ALARM]"],
    ])

    doc.add_heading("7. 黑暗/强光轻量策略", 1)
    doc.add_paragraph("系统从 UYVY 图像中抽样 Y 通道平均亮度 avg_luma。连续达到阈值后进入 LOW_LIGHT 或 BRIGHT 策略。低照下降置信度阈值以减少漏检，强光提高阈值以减少误检。该策略无需额外模型，适合端侧轻量部署。")
    table(["策略", "触发条件", "参数变化"], [
        ["NORMAL", "40 <= avg_luma <= 210", "conf=0.30"],
        ["LOW_LIGHT", "avg_luma < 40 且稳定采样", "conf=0.25"],
        ["BRIGHT", "avg_luma > 210 且稳定采样", "conf=0.35"],
    ])

    doc.add_heading("8. 模型说明与声明", 1)
    doc.add_paragraph("正式模型文件名为 smart_guard_coco_256.m1model。检测模型采用 COCO 通用目标检测权重作为基线，经 Head6 输出适配、M1 模型转换、256x256 输入压缩，并结合 Aurora 现场画面完成阈值、危险区、低照/强光策略和端侧性能验证调优。")
    doc.add_paragraph("当前文档不声称已完成 Aurora 自采样再训练。如果赛题方强制要求自采样训练，应补充 Aurora 采集、标注、微调和对比评测材料。")

    doc.add_heading("9. 遇到的问题与优化方案", 1)
    table(["问题", "定位", "优化"], [
        ["人体框不完整", "YOLOv8 DFL 解码缺少 stride", "恢复 ltrb * stride_f"],
        ["危险区误判", "crop 坐标和 OSD 坐标混用", "判断保持 crop 坐标，显示前 x+240"],
        ["日志/快照影响 FPS", "热路径 I/O 较重", "降低运行态快照频率，后台编码"],
        ["端侧 FPS 偏低", "推理输入和后处理负载偏大", "输入尺寸改为 256"],
        ["现场光照变化", "低照漏检、强光误检", "avg_luma 策略动态调阈值"],
    ])

    doc.add_heading("10. 复现指南", 1)
    for step in [
        "将 smart_guard_coco_256.m1model 放入 /app_demo/app_assets/models/。",
        "编译并部署 ssne_ai_yolo_coco，可使用 ssne_ai_yolo_coco/scripts/run.sh 启动。",
        "PC 端修改 pc_controller/controller_config.json 中的 serial_port。",
        "运行 pc_controller/run.bat，绘制危险区并按 S 发送。",
        "板端进入正式检测后，保存串口日志和 Aurora 调试画面。",
    ]:
        doc.add_paragraph(step)

    doc.add_heading("11. 待补材料清单", 1)
    table(["材料", "用途", "状态"], [
        ["Aurora 调试画面截图", "证明 OSD 和检测效果", "待补"],
        ["串口 60 秒 SUMMARY 截图", "证明性能和稳定性", "待补"],
        ["实物接线图/报警视频帧", "证明 GPIO 声光报警", "待补"],
        ["模型转换平台截图", "证明模型转换流程", "可选补充"],
        ["最终技术数据压缩包目录截图", "证明提交完整性", "待补"],
    ])

    path = OUT / "Smart-Home-Guard_technical_report_draft.docx"
    doc.save(path)
    return path


def build_checklist() -> Path:
    path = OUT / "asset_checklist.md"
    path.write_text(
        "# 待补素材清单\n\n"
        "- Aurora 调试画面截图：正常光、黑暗、强光各 1 张。\n"
        "- 串口 60 秒验收日志截图：包含 [CHECK][SUMMARY] 与 [CHECK][COUNTS]。\n"
        "- 作品演示视频素材：危险区配置、进入危险区、OSD 红框、ALERT、蜂鸣器/LED。\n"
        "- 实物图：A1 板、摄像头、FT232RL 接线、报警外设。\n"
        "- 模型转换证据：smart_guard_coco_head6_256.onnx 转 smart_guard_coco_256.m1model 的平台截图。\n"
        "- 技术数据压缩包截图：源码、模型、可执行文件、脚本、README。\n\n"
        "## 模型声明建议\n\n"
        "检测模型采用 COCO 通用目标检测权重作为基线，经 Head6 输出适配、M1 模型转换、256x256 输入压缩，"
        "并结合 Aurora 现场画面完成阈值、危险区、低照/强光策略和端侧性能验证调优。当前材料不声称已完成 Aurora 自采样再训练。\n",
        encoding="utf-8",
    )
    return path


if __name__ == "__main__":
    for output in [build_pptx(), build_docx(), build_checklist()]:
        print(output)
